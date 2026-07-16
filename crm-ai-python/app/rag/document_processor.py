"""文档处理 — 解析上传文件、文本分块"""

from __future__ import annotations

import os
import re
from pathlib import Path
from typing import BinaryIO


def extract_text_from_file(file_path: str) -> str:
    """根据文件扩展名提取文本内容"""
    ext = Path(file_path).suffix.lower()
    if ext == ".txt":
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    elif ext == ".md":
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    elif ext == ".docx":
        return _extract_docx(file_path)
    elif ext == ".pdf":
        return _extract_pdf(file_path)
    elif ext == ".pptx":
        return _extract_pptx(file_path)
    else:
        raise ValueError(f"不支持的文件格式: {ext}")


def _extract_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs)


def _extract_pdf(path: str) -> str:
    import fitz
    doc = fitz.open(path)
    return "\n".join(page.get_text() for page in doc)


def _extract_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
    return "\n".join(texts)


# ── 文本分块策略 ──

def chunk_text(text: str, chunk_size: int = 512, overlap: int = 64) -> list[str]:
    """
    将长文本按段落和长度分块。
    - 优先按段落（空行）分割
    - 超过 chunk_size 的段落再按句子切割
    - 相邻块间保留 overlap 字符重叠
    """
    if not text.strip():
        return []

    # 按空行分割为段落
    paragraphs = re.split(r"\n\s*\n", text)
    paragraphs = [p.strip() for p in paragraphs if p.strip()]

    chunks: list[str] = []
    buffer = ""

    for para in paragraphs:
        if len(buffer) + len(para) + 1 <= chunk_size:
            buffer = (buffer + "\n" + para).strip()
        else:
            if buffer:
                chunks.append(buffer)
            # 段落本身超长 → 按句切割
            if len(para) > chunk_size:
                sentences = re.split(r"(?<=[。！？.!?])", para)
                sub_buffer = ""
                for sent in sentences:
                    if not sent.strip():
                        continue
                    if len(sub_buffer) + len(sent) + 1 <= chunk_size:
                        sub_buffer += sent
                    else:
                        if sub_buffer:
                            chunks.append(sub_buffer.strip())
                        sub_buffer = sent
                if sub_buffer:
                    buffer = sub_buffer.strip()
                else:
                    buffer = ""
            else:
                buffer = para

    if buffer:
        chunks.append(buffer)

    # 相邻块 overlap
    if overlap > 0 and len(chunks) > 1:
        merged = [chunks[0]]
        for i in range(1, len(chunks)):
            prev = merged[-1]
            curr = chunks[i]
            # 前一块末尾 overlap 字符拼到当前块开头
            if len(prev) > overlap:
                curr = prev[-overlap:] + curr
            merged.append(curr)
        chunks = merged

    return chunks
